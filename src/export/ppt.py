"""PPT 导出：每张图对应一页，背景图 + 文本框。复用 banana-slides 的 PPTXBuilder。"""
import tempfile
from pathlib import Path
from typing import Callable, List, Optional, Union, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

from ..utils.pptx_builder import PPTXBuilder
from .merge_blocks import merge_inline_blocks, merge_vertical_paragraphs


# 默认幻灯片尺寸（英寸，16:9 宽屏），与 banana-slides 一致
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
# 幻灯片逻辑像素（96 DPI），用于 PPTXBuilder
SLIDE_WIDTH_PX = int(SLIDE_WIDTH_IN * 96)
SLIDE_HEIGHT_PX = int(SLIDE_HEIGHT_IN * 96)


def _styled_block_to_text_style(blk: dict) -> object:
    """将本项目的 styled_block 转为 PPTXBuilder 可用的 text_style 对象（兼容 banana-slides TextStyleResult）。"""
    color = blk.get("color", (0, 0, 0))
    if isinstance(color, (list, tuple)) and len(color) >= 3:
        font_color_rgb = (int(color[0]), int(color[1]), int(color[2]))
    else:
        font_color_rgb = (0, 0, 0)
    return type(
        "TextStyle",
        (),
        {
            "font_color_rgb": font_color_rgb,
            "is_bold": bool(blk.get("bold", False)),
            "is_italic": False,
            "is_underline": False,
            "text_alignment": None,
            "colored_segments": None,
        },
    )()


def _px_to_inches(x_px: float, y_px: float, img_w: int, img_h: int, slide_w_in: float, slide_h_in: float) -> tuple:
    """将像素坐标按「图→幻灯片」比例换算为英寸。"""
    scale_x = slide_w_in / max(1, img_w)
    scale_y = slide_h_in / max(1, img_h)
    return (x_px * scale_x, y_px * scale_y)


def _box_bounds(box: List[List[float]]) -> tuple:
    """返回 (x_min, y_min, x_max, y_max)。"""
    xs = [p[0] for p in box]
    ys = [p[1] for p in box]
    return (min(xs), min(ys), max(xs), max(ys))


def build_editable_pptx(
    slides_data: List[Tuple[Image.Image, List[dict], int, int]],
    output_path: Union[str, Path],
    dpi: int = 96,
    progress_callback: Optional[Callable[[str, int, int, str], None]] = None,
    *,
    font_normal: str = "Tencent Sans W3",
    font_bold: str = "Tencent Sans W7",
    font_ea_normal: str = "腾讯字体 W3",
    font_ea_bold: str = "腾讯字体 W7",
    text_lang: str = "zh-CN",
    text_alt_lang: str = "en-US",
    text_pad_ratio: float = 0.08,
    width_safety: float = 0.96,
    merge_textbox: bool = True,
) -> None:
    """
    使用 banana-slides 的 PPTXBuilder 生成可编辑 PPTX。
    每页：干净背景图 + 按 bbox 放置的文本框（可编辑）。
    slides_data: [(background_image, styled_blocks, img_w, img_h), ...]
    progress_callback: 可选，(phase, current, total, message) -> None。

    附加优化：
    - merge_textbox=True 时，先做同行短框合并（同样式相邻，见 merge_blocks）
    - 再按 text_pad_ratio 对单个文本框向右扩宽，避免贴边折行
    - width_safety<1.0 时，反推字号多留 (1-width_safety) 安全余量，
      对抗 OCR bbox 比真实 glyph advance 略紧导致的临界折行
    - 每个 run 同时写入西文/东亚字体与 lang/altLang，避免中英文系统字体回退与拼写误报
    """
    builder = PPTXBuilder()
    builder.setup_presentation_size(SLIDE_WIDTH_PX, SLIDE_HEIGHT_PX, dpi=dpi)
    builder.create_presentation()
    total = len(slides_data)

    for idx, (background_image, styled_blocks, img_w, img_h) in enumerate(slides_data):
        if progress_callback:
            progress_callback("export", idx + 1, total, "写入幻灯片")

        # —— 等比适配（contain）：统一缩放 + 居中偏移，防止背景图非等比拉伸 ——
        # 用同一个 scale 换算背景图与所有文本框，保证二者相对位置不变
        scale = min(
            SLIDE_WIDTH_PX / max(1, img_w),
            SLIDE_HEIGHT_PX / max(1, img_h),
        )
        drawn_w_px = img_w * scale
        drawn_h_px = img_h * scale
        offset_x_px = (SLIDE_WIDTH_PX - drawn_w_px) / 2.0
        offset_y_px = (SLIDE_HEIGHT_PX - drawn_h_px) / 2.0

        slide = builder.add_blank_slide()

        # 背景图：临时保存后按等比缩放尺寸插入并居中，然后移到底层
        if isinstance(background_image, (str, Path)):
            bg_path = str(background_image)
            need_unlink = False
        else:
            fd = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            background_image.save(fd.name, format="PNG")
            bg_path = fd.name
            need_unlink = True
        try:
            # 像素按 dpi 换算成英寸
            left_in = offset_x_px / dpi
            top_in = offset_y_px / dpi
            width_in = drawn_w_px / dpi
            height_in = drawn_h_px / dpi
            pic = slide.shapes.add_picture(
                bg_path,
                Inches(left_in), Inches(top_in),
                width=Inches(width_in),
                height=Inches(height_in),
            )
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(2, pic._element)
        finally:
            if need_unlink:
                Path(bg_path).unlink(missing_ok=True)

        # 1) 合并同行短文本框（在像素坐标上做，阈值按图像分辨率，不受 slide 尺寸影响）
        effective_blocks = styled_blocks
        if merge_textbox and styled_blocks:
            effective_blocks = merge_inline_blocks(
                styled_blocks,
                slide_width_px=img_w,  # 以原图像素为上限，避免跨图边
            )
            # 1b) 再做跨行段落合并：同样式、左对齐、y 紧邻一行内
            effective_blocks = merge_vertical_paragraphs(
                effective_blocks,
                slide_width_px=img_w,
                slide_height_px=img_h,
            )

        # 内容区右下边界（slide 像素坐标系），用于裁切和判定上限
        content_left_px = offset_x_px
        content_top_px = offset_y_px
        content_right_px = offset_x_px + drawn_w_px
        content_bottom_px = offset_y_px + drawn_h_px

        for blk in effective_blocks:
            box = blk.get("box")
            text = (blk.get("text") or "").strip()
            if not box or not text:
                continue
            x0, y0, x1, y1 = _box_bounds(box)
            # 用同一个 scale + 居中偏移换算到 slide 像素坐标
            bx0 = int(x0 * scale + offset_x_px)
            by0 = int(y0 * scale + offset_y_px)
            bx1 = int(x1 * scale + offset_x_px)
            by1 = int(y1 * scale + offset_y_px)
            # 2) 仅向右扩宽，不越内容区右边界（即原图右侧在 slide 上的映射位置）
            if text_pad_ratio and text_pad_ratio > 0:
                pad = int((bx1 - bx0) * text_pad_ratio)
                bx1 = bx1 + pad
            bx1 = min(bx1, int(content_right_px))
            # 纵向也别溢出内容区下边界（极端情况下保个底）
            by1 = min(by1, int(content_bottom_px))
            bbox_slide = [bx0, by0, bx1, by1]

            text_style = _styled_block_to_text_style(blk)
            is_title = bool(blk.get("bold"))
            text_level = "title" if is_title else "default"
            # 3) 根据 bold 选择字体；粗体用 W7，否则 W3
            latin = font_bold if is_title else font_normal
            east_asian = font_ea_bold if is_title else font_ea_normal
            try:
                builder.add_text_element(
                    slide=slide,
                    text=text,
                    bbox=bbox_slide,
                    text_level=text_level,
                    dpi=dpi,
                    align="left",
                    text_style=text_style,
                    font_latin=latin,
                    font_east_asian=east_asian,
                    text_lang=text_lang,
                    text_alt_lang=text_alt_lang,
                    expand_ratio=0.0,  # bbox 外扩已在上方 text_pad_ratio 处理
                    width_safety=width_safety,
                )
            except Exception as e:
                import logging
                logging.getLogger(__name__).warning("Skip text element %s: %s", text[:30], e)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    builder.save(str(output_path))


def add_slide_from_image_and_blocks(
    prs: Presentation,
    background_image: Union[Image.Image, str, Path],
    styled_blocks: List[dict],
    img_w: int,
    img_h: int,
    slide_width_in: float = SLIDE_WIDTH_IN,
    slide_height_in: float = SLIDE_HEIGHT_IN,
    font_normal: str = "Tencent Sans W3",
    font_bold: str = "Tencent Sans W7",
) -> None:
    """
    向 prs 追加一页：以 background_image 为全页底图，再按 styled_blocks 在对应位置画文本框。
    styled_blocks 每项: {"box", "text", "bold", "color", "font_size_pt"}。
    """
    blank = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(blank)

    if isinstance(background_image, (str, Path)):
        img_path = str(background_image)
    else:
        fp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        background_image.save(fp.name, format="PNG")
        img_path = fp.name
    try:
        pic = slide.shapes.add_picture(
            img_path, Inches(0), Inches(0),
            width=Inches(slide_width_in),
            height=Inches(slide_height_in),
        )
        # 将背景图片移到最底层（通过操作 XML 元素）
        spTree = slide.shapes._spTree
        spTree.remove(pic._element)
        spTree.insert(2, pic._element)  # 索引 2：在 nvGrpSpPr 和 grpSpPr 之后
    finally:
        if isinstance(background_image, Image.Image):
            Path(img_path).unlink(missing_ok=True)

    for blk in styled_blocks:
        box = blk.get("box")
        text = blk.get("text", "").strip()
        if not box or not text:
            continue
        bold = bool(blk.get("bold", False))
        color = blk.get("color", (0, 0, 0))
        font_size_pt = float(blk.get("font_size_pt", 18))
        x0, y0, x1, y1 = _box_bounds(box)
        # 文本框左上角与宽高（英寸）
        left_in, top_in = _px_to_inches(x0, y0, img_w, img_h, slide_width_in, slide_height_in)
        w_in = max(0.1, (x1 - x0) / max(1, img_w) * slide_width_in)
        h_in = max(0.1, (y1 - y0) / max(1, img_h) * slide_height_in)
        tb = slide.shapes.add_textbox(
            Inches(left_in), Inches(top_in),
            Inches(w_in), Inches(h_in),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold
        run.font.name = font_bold if bold else font_normal
        if isinstance(color, (list, tuple)) and len(color) >= 3:
            run.font.color.rgb = RGBColor(color[0], color[1], color[2])
