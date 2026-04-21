"""
PPTX Builder - 创建可编辑 PPTX
复用自 banana-slides (Based on OpenDCAI/DataFlow-Agent)
"""
import os
import logging
from datetime import datetime, timezone
from typing import List, Dict, Any, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from html.parser import HTMLParser

logger = logging.getLogger(__name__)

# 可选：PIL ImageFont 仅用于精确测量文字宽度，缺失时用估算
try:
    from PIL import Image, ImageFont
    _HAS_PIL_FONT = True
except ImportError:
    _HAS_PIL_FONT = False


class HTMLTableParser(HTMLParser):
    """Parse HTML table into row/column data"""

    def __init__(self):
        super().__init__()
        self.table_data = []
        self.current_row = []
        self.current_cell = []
        self.in_table = False
        self.in_row = False
        self.in_cell = False

    def handle_starttag(self, tag, attrs):
        if tag == "table":
            self.in_table = True
            self.table_data = []
        elif tag == "tr":
            self.in_row = True
            self.current_row = []
        elif tag in ["td", "th"]:
            self.in_cell = True
            self.current_cell = []

    def handle_endtag(self, tag):
        if tag == "table":
            self.in_table = False
        elif tag == "tr":
            self.in_row = False
            if self.current_row:
                self.table_data.append(self.current_row)
        elif tag in ["td", "th"]:
            self.in_cell = False
            cell_text = "".join(self.current_cell).strip()
            self.current_row.append(cell_text)

    def handle_data(self, data):
        if self.in_cell:
            self.current_cell.append(data)

    @staticmethod
    def parse_html_table(html: str) -> List[List[str]]:
        """Parse HTML table string into 2D array of cells"""
        parser = HTMLTableParser()
        parser.feed(html)
        return parser.table_data


class PPTXBuilder:
    """Builder for creating editable PPTX from structured content (from banana-slides)"""

    DEFAULT_SLIDE_WIDTH_INCHES = 10
    DEFAULT_SLIDE_HEIGHT_INCHES = 5.625
    DEFAULT_DPI = 96
    MAX_SLIDE_WIDTH_INCHES = 56.0
    MAX_SLIDE_HEIGHT_INCHES = 56.0
    MIN_SLIDE_WIDTH_INCHES = 1.0
    MIN_SLIDE_HEIGHT_INCHES = 1.0
    MIN_FONT_SIZE = 6
    MAX_FONT_SIZE = 200

    # 可选字体路径（不存在时用字符估算）
    FONT_PATH = os.path.join(os.path.dirname(__file__), "..", "fonts", "NotoSansSC-Regular.ttf")
    _font_cache: Dict[float, Any] = {}

    @classmethod
    def _get_font(cls, size_pt: float) -> Optional[Any]:
        if not _HAS_PIL_FONT or not os.path.exists(cls.FONT_PATH):
            return None
        cache_key = round(size_pt * 2) / 2
        if cache_key not in cls._font_cache:
            try:
                cls._font_cache[cache_key] = ImageFont.truetype(cls.FONT_PATH, int(size_pt))
            except Exception as e:
                logger.warning(f"Failed to load font {cls.FONT_PATH}: {e}")
                return None
        return cls._font_cache[cache_key]

    @classmethod
    def _measure_text_width(cls, text: str, font_size_pt: float) -> Optional[float]:
        font = cls._get_font(font_size_pt)
        if font is None:
            return None
        try:
            bbox = font.getbbox(text)
            return bbox[2] - bbox[0]
        except Exception as e:
            logger.warning(f"Failed to measure text: {e}")
            return None

    def __init__(self, slide_width_inches: float = None, slide_height_inches: float = None):
        self.slide_width_inches = slide_width_inches or self.DEFAULT_SLIDE_WIDTH_INCHES
        self.slide_height_inches = slide_height_inches or self.DEFAULT_SLIDE_HEIGHT_INCHES
        self.prs = None
        self.current_slide = None

    def create_presentation(self) -> Presentation:
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width_inches)
        self.prs.slide_height = Inches(self.slide_height_inches)
        self._set_core_properties(self.prs)
        return self.prs

    @staticmethod
    def _set_core_properties(prs: Presentation) -> None:
        try:
            core = prs.core_properties
            now = datetime.now(timezone.utc)
            core.author = "images-2-ppt"
            core.last_modified_by = "images-2-ppt"
            core.created = now
            core.modified = now
            # 不要给 last_printed 赋 None：python-pptx 会校验 datetime 类型。
            # 保持默认（未设置）即可。
        except Exception as e:
            logger.warning(f"Failed to set core properties: {e}")

    def setup_presentation_size(self, width_pixels: int, height_pixels: int, dpi: int = None):
        dpi = dpi or self.DEFAULT_DPI
        width_inches = width_pixels / dpi
        height_inches = height_pixels / dpi
        scale_factor = 1.0
        if width_inches > self.MAX_SLIDE_WIDTH_INCHES:
            scale_factor = self.MAX_SLIDE_WIDTH_INCHES / width_inches
            logger.warning(f"Slide width exceeds limit, scaling by {scale_factor:.3f}x")
        height_scale = 1.0
        if height_inches > self.MAX_SLIDE_HEIGHT_INCHES:
            height_scale = self.MAX_SLIDE_HEIGHT_INCHES / height_inches
            if height_scale < scale_factor:
                scale_factor = height_scale
            logger.warning(f"Slide height exceeds limit, scaling by {scale_factor:.3f}x")
        if scale_factor < 1.0:
            width_inches *= scale_factor
            height_inches *= scale_factor
        width_inches = max(self.MIN_SLIDE_WIDTH_INCHES, width_inches)
        height_inches = max(self.MIN_SLIDE_HEIGHT_INCHES, height_inches)
        self.slide_width_inches = width_inches
        self.slide_height_inches = height_inches
        if self.prs:
            self.prs.slide_width = Inches(self.slide_width_inches)
            self.prs.slide_height = Inches(self.slide_height_inches)

    def add_blank_slide(self):
        if not self.prs:
            self.create_presentation()
        blank_layout = self.prs.slide_layouts[6]
        self.current_slide = self.prs.slides.add_slide(blank_layout)
        return self.current_slide

    def pixels_to_inches(self, pixels: float, dpi: int = None) -> float:
        dpi = dpi or self.DEFAULT_DPI
        return pixels / dpi

    def calculate_font_size(
        self, bbox: List[int], text: str, text_level: Any = None, dpi: int = None,
        width_safety: float = 0.96,
    ) -> float:
        """根据 bbox 和文本反推字号。

        width_safety: 可用宽度系数（<=1.0）。用于在 OCR bbox 宽度比真实
        glyph advance 略紧的情况下，给字号再打点保险，避免临界折行。
        默认 0.96（约 4% 余量），极端缩窄可调到 0.90。
        """
        dpi = dpi or self.DEFAULT_DPI
        width_px = bbox[2] - bbox[0]
        height_px = bbox[3] - bbox[1]
        width_pt = (width_px / dpi) * 72
        height_pt = (height_px / dpi) * 72
        # 宽度留出安全余量，防止 OCR bbox 比真实 glyph advance 略紧导致折行
        safety = max(0.5, min(1.0, float(width_safety)))
        usable_width_pt = width_pt * safety
        usable_height_pt = height_pt
        if usable_width_pt <= 0 or usable_height_pt <= 0:
            return self.MIN_FONT_SIZE
        use_precise = os.path.exists(self.FONT_PATH) and _HAS_PIL_FONT
        best_size = self.MIN_FONT_SIZE
        for font_size in range(int(self.MAX_FONT_SIZE), int(self.MIN_FONT_SIZE) - 1, -1):
            font_size = float(font_size)
            lines = text.split("\n")
            total_required_lines = 0
            for line in lines:
                if not line:
                    total_required_lines += 1
                    continue
                if use_precise:
                    line_width_pt = self._measure_text_width(line, font_size)
                    if line_width_pt is None:
                        use_precise = False
                if not use_precise:
                    cjk_count = sum(
                        1
                        for c in line
                        if "\u4e00" <= c <= "\u9fff"
                        or "\u3040" <= c <= "\u30ff"
                        or "\uac00" <= c <= "\ud7af"
                    )
                    non_cjk_count = len(line) - cjk_count
                    line_width_pt = (cjk_count * 1.0 + non_cjk_count * 0.5) * font_size
                lines_needed = max(1, -(-int(line_width_pt) // int(usable_width_pt)))
                total_required_lines += lines_needed
            line_height_pt = font_size * 1.0
            total_height_pt = total_required_lines * line_height_pt
            if total_height_pt <= usable_height_pt:
                best_size = font_size
                break
        return best_size

    @staticmethod
    def _decorate_run_rpr(
        run,
        *,
        font_latin: Optional[str] = None,
        font_east_asian: Optional[str] = None,
        text_lang: Optional[str] = None,
        text_alt_lang: Optional[str] = None,
    ) -> None:
        """在 run 的 a:rPr 上写入东亚字体与语言标签。

        - a:latin / a:ea 决定英文与中文字形；
        - rPr@lang / rPr@altLang 决定拼写检查使用的语言。
        旧节点会先移除，避免重复。
        """
        # python-pptx 在设置 font 相关属性时才会创建 rPr
        rPr = run._r.get_or_add_rPr()
        if font_latin:
            for old in rPr.findall(qn("a:latin")):
                rPr.remove(old)
            latin = etree.SubElement(rPr, qn("a:latin"))
            latin.set("typeface", font_latin)
        if font_east_asian:
            for old in rPr.findall(qn("a:ea")):
                rPr.remove(old)
            ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", font_east_asian)
        if text_lang:
            rPr.set("lang", text_lang)
        if text_alt_lang:
            rPr.set("altLang", text_alt_lang)

    def add_text_element(
        self,
        slide,
        text: str,
        bbox: List[int],
        text_level: Any = None,
        dpi: int = None,
        align: str = "left",
        text_style: Any = None,
        font_latin: Optional[str] = None,
        font_east_asian: Optional[str] = None,
        text_lang: Optional[str] = "zh-CN",
        text_alt_lang: Optional[str] = "en-US",
        expand_ratio: float = 0.0,
        width_safety: float = 0.96,
    ):
        dpi = dpi or self.DEFAULT_DPI
        has_colored_segments = (
            text_style
            and hasattr(text_style, "colored_segments")
            and text_style.colored_segments
        )
        actual_text = (
            "".join(seg.text for seg in text_style.colored_segments)
            if has_colored_segments
            else text
        )
        # 文本框尺寸：仅按 expand_ratio 轻微外扩（默认 0，由调用侧自己控制外扩）
        EXPAND_RATIO = max(0.0, float(expand_ratio))
        bbox_width = bbox[2] - bbox[0]
        bbox_height = bbox[3] - bbox[1]
        expand_w = bbox_width * EXPAND_RATIO
        expand_h = bbox_height * EXPAND_RATIO
        left = Inches(self.pixels_to_inches(bbox[0] - expand_w / 2, dpi))
        top = Inches(self.pixels_to_inches(bbox[1] - expand_h / 2, dpi))
        width = Inches(self.pixels_to_inches(bbox_width + expand_w, dpi))
        height = Inches(self.pixels_to_inches(bbox_height + expand_h, dpi))
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)

        def replace_some_chars(s: str) -> str:
            return s.replace("·", "•", 1) if s.lstrip().startswith("·") else s

        actual_text = replace_some_chars(actual_text)
        font_size = self.calculate_font_size(bbox, actual_text, text_level, dpi, width_safety=width_safety)
        effective_align = align
        if text_style and getattr(text_style, "text_alignment", None):
            effective_align = text_style.text_alignment
        is_bold = getattr(text_style, "is_bold", False) if text_style else False
        is_italic = getattr(text_style, "is_italic", False) if text_style else False
        is_underline = getattr(text_style, "is_underline", False) if text_style else False
        if text_level == 1 or text_level == "title":
            is_bold = True

        # 决定最终使用的字体（粗体走 bold 字体）
        chosen_latin = font_latin
        chosen_ea = font_east_asian

        if has_colored_segments:
            paragraph = text_frame.paragraphs[0]
            paragraph.clear()
            for seg in text_style.colored_segments:
                run = paragraph.add_run()
                run.text = replace_some_chars(seg.text)
                run.font.size = Pt(font_size)
                run.font.bold = is_bold
                run.font.underline = is_underline
                r, g, b = seg.color_rgb
                run.font.color.rgb = RGBColor(r, g, b)
                run.font.italic = getattr(seg, "is_latex", False) or is_italic
                if chosen_latin:
                    run.font.name = chosen_latin
                self._decorate_run_rpr(
                    run,
                    font_latin=chosen_latin,
                    font_east_asian=chosen_ea,
                    text_lang=text_lang,
                    text_alt_lang=text_alt_lang,
                )
        else:
            # 手动按 "\n" 拆段，保证每行一个段落，对齐/字号等属性逐段设置
            lines = actual_text.split("\n") if actual_text else [""]
            # 清空默认首段
            text_frame.clear()
            first_para = text_frame.paragraphs[0]
            paragraphs_to_style: List[Any] = []
            for idx, line in enumerate(lines):
                p = first_para if idx == 0 else text_frame.add_paragraph()
                paragraphs_to_style.append(p)
                run = p.add_run()
                run.text = line
                run.font.size = Pt(font_size)
                run.font.bold = is_bold
                run.font.italic = is_italic
                run.font.underline = is_underline
                if text_style and getattr(text_style, "font_color_rgb", None):
                    r, g, b = text_style.font_color_rgb
                    run.font.color.rgb = RGBColor(r, g, b)
                if chosen_latin:
                    run.font.name = chosen_latin
                self._decorate_run_rpr(
                    run,
                    font_latin=chosen_latin,
                    font_east_asian=chosen_ea,
                    text_lang=text_lang,
                    text_alt_lang=text_alt_lang,
                )
            # 供后续 align 设置（沿用老接口的 `paragraph` 变量名，指向首段即可）
            paragraph = first_para

        if effective_align == "center":
            align_val = PP_ALIGN.CENTER
        elif effective_align == "right":
            align_val = PP_ALIGN.RIGHT
        elif effective_align == "justify":
            align_val = PP_ALIGN.JUSTIFY
        else:
            align_val = PP_ALIGN.LEFT
        # 多段：每段都要应用对齐
        for p in text_frame.paragraphs:
            p.alignment = align_val

    def add_image_element(self, slide, image_path: str, bbox: List[int], dpi: int = None):
        dpi = dpi or self.DEFAULT_DPI
        if not os.path.exists(image_path):
            self.add_image_placeholder(slide, bbox, dpi)
            return
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            logger.error(f"Failed to add image {image_path}: {e}")
            self.add_image_placeholder(slide, bbox, dpi)

    def add_image_placeholder(self, slide, bbox: List[int], dpi: int = None):
        dpi = dpi or self.DEFAULT_DPI
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = "[Image]"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.italic = True

    def add_table_element(
        self, slide, html_table: str, bbox: List[int], dpi: int = None
    ):
        dpi = dpi or self.DEFAULT_DPI
        try:
            table_data = HTMLTableParser.parse_html_table(html_table)
        except Exception as e:
            logger.error(f"Failed to parse HTML table: {e}")
            return
        if not table_data or not table_data[0]:
            return
        rows, cols = len(table_data), len(table_data[0])
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        try:
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < cols:
                        table.cell(row_idx, col_idx).text = cell_text
        except Exception as e:
            logger.error(f"Failed to create table: {e}")

    def save(self, output_path: str):
        if not self.prs:
            raise ValueError("No presentation to save")
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)

    def get_presentation(self) -> Presentation:
        return self.prs
